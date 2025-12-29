"""
File Parser Module
Extracts text content from various file formats
"""

import os
import sys
from pathlib import Path

# Try to import required libraries
try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
    import PyPDF2
except ImportError as e:
    print(f"Error importing libraries: {e}")
    print("Make sure to install: python-docx, openpyxl, python-pptx, PyPDF2")
    sys.exit(1)


def extract_text_from_pdf(file_path):
    """Extract text from PDF files"""
    try:
        text = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return " ".join(text)
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return ""


def extract_text_from_docx(file_path):
    """Extract text from DOCX files"""
    try:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return " ".join(text)
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
        return ""


def extract_text_from_xlsx(file_path):
    """Extract text from XLSX files"""
    try:
        wb = load_workbook(file_path)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell:
                        text.append(str(cell))
        return " ".join(text)
    except Exception as e:
        print(f"Error reading XLSX {file_path}: {e}")
        return ""


def extract_text_from_pptx(file_path):
    """Extract text from PPTX files"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return " ".join(text)
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        return ""


def extract_text_from_markdown(file_path):
    """Extract text from Markdown files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error reading Markdown {file_path}: {e}")
        return ""


def extract_text_from_text(file_path):
    """Extract text from TXT files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error reading TXT {file_path}: {e}")
        return ""


def extract_content(file_path):
    """
    Main function to extract content from any supported file format
    Returns: (filename, content_text)
    """
    file_path = Path(file_path)
    
    if not file_path.exists():
        print(f"⚠️  File not found: {file_path}")
        return None, None
    
    filename = file_path.name
    extension = file_path.suffix.lower()
    
    # Extract content based on file type
    if extension == ".pdf":
        content = extract_text_from_pdf(file_path)
    elif extension == ".docx":
        content = extract_text_from_docx(file_path)
    elif extension == ".xlsx":
        content = extract_text_from_xlsx(file_path)
    elif extension == ".pptx":
        content = extract_text_from_pptx(file_path)
    elif extension == ".md":
        content = extract_text_from_markdown(file_path)
    elif extension == ".txt":
        content = extract_text_from_text(file_path)
    else:
        print(f"[!] Unsupported file format: {extension}")
        return filename, ""
    
    return filename, content


if __name__ == "__main__":
    # Test the parser
    test_file = Path.home() / "Desktop" / "Semester_Transcript_2024.pdf"
    if test_file.exists():
        print(f"Testing parser with: {test_file.name}")
        name, content = extract_content(test_file)
        print(f"Filename: {name}")
        print(f"Content length: {len(content)} characters")
        print(f"Content preview: {content[:200]}...")
    else:
        print(f"Test file not found: {test_file}")
