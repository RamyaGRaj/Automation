"""
Script to create test dataset in proper file formats on Desktop
"""
import os
from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# Get Desktop path
desktop_path = Path.home() / "Desktop"

print("📁 Creating files on Desktop...")
print(f"Path: {desktop_path}\n")

# ============ UNIVERSITY DOCS ============

# 1. Semester_Transcript_2024.pdf
print("Creating Semester_Transcript_2024.pdf...")
pdf_path = desktop_path / "Semester_Transcript_2024.pdf"
c = canvas.Canvas(str(pdf_path), pagesize=letter)
c.setFont("Helvetica-Bold", 14)
c.drawString(50, 750, "Academic Transcript - 2024")
c.setFont("Helvetica", 11)
c.drawString(50, 720, "Student ID: 12345678")
c.drawString(50, 700, "Fall Semester 2024 GPA: 3.8")
c.drawString(50, 680, "Courses Completed:")
c.drawString(70, 660, "• Data Structures (A)")
c.drawString(70, 640, "• Algorithms (A)")
c.drawString(70, 620, "• Database Systems (A-)")
c.drawString(70, 600, "• Web Development (B+)")
c.drawString(50, 570, "Academic Standing: Excellent")
c.drawString(50, 550, "Cumulative GPA: 3.85")
c.save()
print(f"✅ Created: {pdf_path.name}")

# 2. Course_Registration_Form.docx
print("Creating Course_Registration_Form.docx...")
doc = Document()
doc.add_heading("Spring 2025 Course Registration Form", 0)
doc.add_paragraph("Student ID: 12345678")
doc.add_paragraph("Name: Student Name")
doc.add_heading("Registered Courses:", level=2)
doc.add_paragraph("Computer Networks (CS 301)", style='List Bullet')
doc.add_paragraph("Web Development (CS 305)", style='List Bullet')
doc.add_paragraph("University Physics II (PHYS 202)", style='List Bullet')
doc.add_paragraph("Data Science Fundamentals (CS 310)", style='List Bullet')
doc.add_paragraph("Registration Date: 2025-01-15")
doc.add_paragraph("Status: APPROVED")
doc.save(str(desktop_path / "Course_Registration_Form.docx"))
print(f"✅ Created: Course_Registration_Form.docx")

# 3. Internship_Approval_Letter.pdf
print("Creating Internship_Approval_Letter.pdf...")
pdf_path = desktop_path / "Internship_Approval_Letter.pdf"
c = canvas.Canvas(str(pdf_path), pagesize=letter)
c.setFont("Helvetica-Bold", 12)
c.drawString(50, 750, "INTERNSHIP APPROVAL LETTER")
c.setFont("Helvetica", 11)
c.drawString(50, 720, "Dear Student,")
c.drawString(50, 690, "This letter certifies that your internship application has been APPROVED.")
c.drawString(50, 660, "Company: TechCorp Inc.")
c.drawString(50, 640, "Position: Software Development Intern")
c.drawString(50, 620, "Duration: 3 months (January - March 2025)")
c.drawString(50, 600, "Start Date: January 6, 2025")
c.drawString(50, 580, "Supervisor: John Smith")
c.drawString(50, 550, "You are authorized to begin your internship program immediately.")
c.drawString(50, 520, "Sincerely,")
c.drawString(50, 500, "Internship Coordinator")
c.save()
print(f"✅ Created: Internship_Approval_Letter.pdf")

# 4. Student_ID_Application.docx
print("Creating Student_ID_Application.docx...")
doc = Document()
doc.add_heading("Student ID Application Form", 0)
doc.add_heading("University Student Identification Request", level=2)
doc.add_paragraph("Full Name: Student Name")
doc.add_paragraph("Date of Birth: 01/01/2000")
doc.add_paragraph("Student ID Number: 12345678")
doc.add_paragraph("Academic Level: Senior")
doc.add_paragraph("Major: Computer Science")
doc.add_paragraph("Certification of Enrollment: CONFIRMED")
doc.add_paragraph("Application Status: APPROVED")
doc.add_paragraph("Issued: December 2024")
doc.save(str(desktop_path / "Student_ID_Application.docx"))
print(f"✅ Created: Student_ID_Application.docx")

# ============ TECHNICAL WORK ============

# 5. API_Documentation_Guide.md
print("Creating API_Documentation_Guide.md...")
md_content = """# API Documentation Guide

## RESTful API Reference v1.0

**Base URL:** https://api.example.com/v1

### Authentication
- Bearer Token required
- Include in Authorization header

### Endpoints

#### 1. GET /users
- Description: Retrieve all users
- Response: 200 OK

#### 2. POST /data/process
- Description: Process technical data
- Request body: JSON
- Content-Type: application/json

#### 3. Docker Configuration
- Container management
- Image building
- Deployment automation

### Error Handling
- 400: Bad Request
- 401: Unauthorized
- 404: Not Found
- 500: Internal Server Error
"""
with open(str(desktop_path / "API_Documentation_Guide.md"), "w") as f:
    f.write(md_content)
print(f"✅ Created: API_Documentation_Guide.md")

# 6. DevOps_Automation_Notes.txt
print("Creating DevOps_Automation_Notes.txt...")
txt_content = """DevOps Automation Notes - December 2025

CI/CD Pipeline Configuration:
- GitHub Actions workflow setup
- Automated testing on every commit
- Docker container building
- Infrastructure deployment to AWS

Docker Container Orchestration:
- Container networking
- Volume management
- Service scaling
- Production deployment

Key Technologies:
- Docker
- Kubernetes
- Jenkins
- GitHub Actions
- Terraform (Infrastructure as Code)

Automation Benefits:
- Reduced manual deployment errors
- Faster release cycles
- Improved system reliability
- Better resource utilization
"""
with open(str(desktop_path / "DevOps_Automation_Notes.txt"), "w") as f:
    f.write(txt_content)
print(f"✅ Created: DevOps_Automation_Notes.txt")

# 7. Docker_Configuration_Checklist.docx
print("Creating Docker_Configuration_Checklist.docx...")
doc = Document()
doc.add_heading("Docker Configuration Checklist", 0)
doc.add_heading("Production Ready Docker Setup", level=2)
doc.add_paragraph("Dockerfile created with proper base image", style='List Bullet')
doc.add_paragraph(".dockerignore file configured", style='List Bullet')
doc.add_paragraph("Image built and tested locally", style='List Bullet')
doc.add_paragraph("Container networking configured", style='List Bullet')
doc.add_paragraph("Volume mounts for persistent data", style='List Bullet')
doc.add_paragraph("Environment variables set", style='List Bullet')
doc.add_paragraph("Security: Run as non-root user", style='List Bullet')
doc.add_paragraph("Resource limits: CPU and memory", style='List Bullet')
doc.add_paragraph("Logging configuration", style='List Bullet')
doc.add_paragraph("Health checks implemented", style='List Bullet')
doc.add_paragraph("Deployment to registry", style='List Bullet')
doc.add_paragraph("Docker Compose for multi-container setup", style='List Bullet')
doc.add_paragraph("CI/CD integration", style='List Bullet')
doc.add_paragraph("Monitoring and alerting", style='List Bullet')
doc.add_paragraph()
doc.add_paragraph("Status: Configuration Complete")
doc.add_paragraph("Date: December 2025")
doc.save(str(desktop_path / "Docker_Configuration_Checklist.docx"))
print(f"✅ Created: Docker_Configuration_Checklist.docx")

# ============ CAPSTONE WORK ============

# 8. Capstone_Project_Proposal.pdf
print("Creating Capstone_Project_Proposal.pdf...")
pdf_path = desktop_path / "Capstone_Project_Proposal.pdf"
c = canvas.Canvas(str(pdf_path), pagesize=letter)
c.setFont("Helvetica-Bold", 14)
c.drawString(50, 750, "Capstone Project Proposal")
c.setFont("Helvetica-Bold", 12)
c.drawString(50, 720, "Title: Intelligent Document Classification System Using Machine Learning")
c.setFont("Helvetica-Bold", 11)
c.drawString(50, 690, "Abstract:")
c.setFont("Helvetica", 10)
c.drawString(70, 670, "This research proposes an automated document classification system")
c.drawString(70, 655, "utilizing semantic analysis and machine learning to categorize")
c.drawString(70, 640, "documents with high accuracy.")
c.setFont("Helvetica-Bold", 11)
c.drawString(50, 610, "Objectives:")
c.setFont("Helvetica", 10)
c.drawString(70, 590, "1. Build a semantic keyword-based classifier")
c.drawString(70, 575, "2. Achieve 95%+ classification accuracy")
c.drawString(70, 560, "3. Support multiple document formats")
c.drawString(70, 545, "4. Enable cross-platform automation")
c.setFont("Helvetica-Bold", 11)
c.drawString(50, 515, "Research Methodology:")
c.setFont("Helvetica", 10)
c.drawString(70, 495, "- Dataset: 500 documents across 3 categories")
c.drawString(70, 480, "- Features: Filename analysis + content extraction")
c.drawString(70, 465, "- Algorithm: Keyword matching with semantic scoring")
c.drawString(70, 450, "- Validation: 5-run consistency testing")
c.save()
print(f"✅ Created: Capstone_Project_Proposal.pdf")

# 9. Capstone_Data_Collection_Log.xlsx
print("Creating Capstone_Data_Collection_Log.xlsx...")
wb = Workbook()
ws = wb.active
ws.title = "Data Collection"
ws.append(["Project", "Intelligent Document Classification"])
ws.append(["Collection Period", "December 2025"])
ws.append([])
ws.append(["Date", "Samples", "Total"])
ws.append(["2025-12-15", 100, 100])
ws.append(["2025-12-20", 100, 200])
ws.append(["2025-12-25", 150, 350])
ws.append(["2025-12-29", 150, 500])
ws.append([])
ws.append(["Category", "Count"])
ws.append(["University Documents", 150])
ws.append(["Technical Documents", 200])
ws.append(["Capstone Documents", 150])
ws.append(["Total", 500])
wb.save(str(desktop_path / "Capstone_Data_Collection_Log.xlsx"))
print(f"✅ Created: Capstone_Data_Collection_Log.xlsx")

# 10. Final_Presentation_Slides.pptx
print("Creating Final_Presentation_Slides.pptx...")
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Project Overview
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
title1 = slide1.shapes.title
content1 = slide1.placeholders[1]
title1.text = "Capstone Project Overview"
content1.text = "Intelligent Document Classification System\n\nObjective: Automated categorization of documents\n\nDuration: 10 minutes presentation + 5 minutes Q&A"

# Slide 2: Research Methodology
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]
title2.text = "Research Methodology"
content2.text = "• Semantic keyword-based approach\n• Content extraction from multiple file formats\n• Confidence scoring algorithm\n• 5-run validation for consistency"

# Slide 3: Key Findings
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]
title3.text = "Key Findings"
content3.text = "✓ 100% accuracy achieved on test dataset\n✓ Consistent performance across 5 runs\n✓ Support for PDF, DOCX, XLSX, PPTX, MD\n✓ Real-world applicability confirmed"

# Slide 4: Results & Impact
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]
title4.text = "Results & Impact"
content4.text = "✓ Desktop automation successful\n✓ Google Drive integration planned\n✓ GitHub repository created\n✓ Production-ready system delivered"

# Slide 5: Conclusion
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
content5 = slide5.placeholders[1]
title5.text = "Conclusion"
content5.text = "✓ Project successfully completed\n✓ All objectives met\n✓ Future enhancements identified\n✓ Ready for deployment"

prs.save(str(desktop_path / "Final_Presentation_Slides.pptx"))
print(f"✅ Created: Final_Presentation_Slides.pptx")

# ============ SUMMARY ============
print("\n" + "="*60)
print("✅ ALL FILES CREATED SUCCESSFULLY!")
print("="*60)
print("\n📁 UNIVERSITY DOCS (4 files):")
print("  ✅ Semester_Transcript_2024.pdf")
print("  ✅ Course_Registration_Form.docx")
print("  ✅ Internship_Approval_Letter.pdf")
print("  ✅ Student_ID_Application.docx")
print("\n📁 TECHNICAL WORK (3 files):")
print("  ✅ API_Documentation_Guide.md")
print("  ✅ DevOps_Automation_Notes.txt")
print("  ✅ Docker_Configuration_Checklist.docx")
print("\n📁 CAPSTONE WORK (3 files):")
print("  ✅ Capstone_Project_Proposal.pdf")
print("  ✅ Capstone_Data_Collection_Log.xlsx")
print("  ✅ Final_Presentation_Slides.pptx")
print(f"\nLocation: {desktop_path}")
print("="*60)
