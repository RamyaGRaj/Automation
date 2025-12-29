"""
Dataset Setup and Verification Script
Ensures all 10 test files are on Desktop root (not in folders)
"""

import os
import shutil
from pathlib import Path
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

DESKTOP_PATH = os.path.expanduser("~/Desktop")

# Create empty folders
FOLDERS = ["UNIVERSITY_DOCS", "TECHNICAL_WORK", "CAPSTONE_WORK"]
for folder in FOLDERS:
    folder_path = os.path.join(DESKTOP_PATH, folder)
    os.makedirs(folder_path, exist_ok=True)

# Move any files from folders back to root
print("🔄 Resetting dataset...")
for folder in FOLDERS:
    folder_path = os.path.join(DESKTOP_PATH, folder)
    if os.path.isdir(folder_path):
        for file in os.listdir(folder_path):
            src = os.path.join(folder_path, file)
            dst = os.path.join(DESKTOP_PATH, file)
            if os.path.isfile(src):
                try:
                    shutil.move(src, dst)
                    print(f"   ↩️  Moved {file} to Desktop root")
                except Exception as e:
                    print(f"   ⚠️  Could not move {file}: {e}")

# Files to create with their content
files_to_create = {
    "Semester_Transcript_2024.pdf": ("UNIVERSITY_DOCS", """
    UNIVERSITY OF TECHNOLOGY
    OFFICIAL SEMESTER TRANSCRIPT
    Student Name: John Smith
    Student ID: 2024-001
    Semester: Fall 2024
    
    COURSES:
    - CS101 Introduction to Computer Science: A
    - MTH201 Calculus II: A-
    - PHY101 Physics I: B+
    - ENG101 English Composition: A
    
    GPA: 3.75/4.0
    Status: Good Standing
    Date Issued: December 29, 2024
    """),
    
    "Course_Registration_Form.docx": ("UNIVERSITY_DOCS", """
    COURSE REGISTRATION FORM
    University of Technology - Spring 2025
    
    Student Information:
    Name: Sarah Johnson
    Student ID: 2024-002
    Program: Computer Science
    
    Courses Registering For:
    1. CS202 Data Structures
    2. MTH202 Linear Algebra
    3. CS203 Web Development
    4. PHY202 Physics II
    
    Advisor Signature: _______________
    Date: December 29, 2024
    """),
    
    "Internship_Approval_Letter.pdf": ("UNIVERSITY_DOCS", """
    INTERNSHIP APPROVAL LETTER
    
    To Whom It May Concern:
    
    This is to certify that Michael Chen, Student ID 2024-003, has been approved
    for a summer internship program at TechCorp Inc.
    
    Position: Software Engineering Intern
    Duration: June 2025 - August 2025
    Supervisor: Dr. Patricia Williams
    
    The student has completed all prerequisite courses and has demonstrated
    excellent academic performance. We approve this internship as part of
    the degree requirements.
    
    Signed,
    Dean of Engineering
    University of Technology
    December 29, 2024
    """),
    
    "Student_ID_Application.docx": ("UNIVERSITY_DOCS", """
    STUDENT ID APPLICATION FORM
    
    Applicant Information:
    Full Name: Emma Wilson
    Date of Birth: January 15, 2003
    Program: Business Administration
    Entry Year: 2021
    Current Year: Third Year
    
    Contact:
    Email: emma.wilson@university.edu
    Phone: +1-555-0123
    
    Declaration:
    I hereby declare that the information provided is accurate and complete.
    
    Signature: ________________
    Date: December 29, 2024
    """),
    
    "API_Documentation_Guide.md": ("TECHNICAL_WORK", """
    # REST API Documentation Guide
    
    ## Overview
    Comprehensive guide for building and documenting RESTful APIs.
    
    ## Endpoints
    
    ### GET /api/users
    Retrieve list of all users.
    
    ### POST /api/users
    Create new user.
    
    ### GET /api/users/{id}
    Retrieve user by ID.
    
    ## Authentication
    Use Bearer token in Authorization header:
    Authorization: Bearer {token}
    
    ## Response Format
    All responses return JSON format.
    Success: {"status": 200, "data": {...}}
    Error: {"status": 400, "error": "Bad Request"}
    
    ## Rate Limiting
    API is rate limited to 1000 requests per hour.
    """),
    
    "DevOps_Automation_Notes.txt": ("TECHNICAL_WORK", """
    DevOps Automation Notes - December 2024
    
    Infrastructure Setup:
    - Configured Docker containers for microservices
    - Set up CI/CD pipeline with Jenkins
    - Automated database backups using cron jobs
    - Implemented monitoring with Prometheus and Grafana
    
    Deployment Process:
    1. Code commit triggers automated tests
    2. Tests pass -> build Docker image
    3. Push to container registry
    4. Deploy to staging environment
    5. Run integration tests
    6. Deploy to production
    
    Rollback Procedure:
    - Keep previous image versions in registry
    - Use kubectl rollout undo for Kubernetes deployments
    - Monitor metrics during rollback
    
    Lessons Learned:
    - Always test rollback procedures
    - Automate as much as possible
    - Monitor everything in production
    """),
    
    "Docker_Configuration_Checklist.docx": ("TECHNICAL_WORK", """
    DOCKER CONFIGURATION CHECKLIST
    
    Pre-Deployment Verification:
    ☑ Dockerfile created and tested locally
    ☑ Base image is security patched
    ☑ Environment variables configured
    ☑ Health checks defined
    ☑ Resource limits set (CPU, memory)
    ☑ Volume mounts configured correctly
    
    Registry Configuration:
    ☑ Image tagged with version number
    ☑ Image pushed to container registry
    ☑ Image signed for security
    
    Orchestration:
    ☑ Kubernetes manifests created
    ☑ Service endpoints configured
    ☑ Ingress rules defined
    ☑ Persistent volumes mounted
    
    Monitoring:
    ☑ Logging configured
    ☑ Metrics collection enabled
    ☑ Alerting rules set
    
    Completed: December 29, 2024
    """),
    
    "Capstone_Project_Proposal.pdf": ("CAPSTONE_WORK", """
    CAPSTONE PROJECT PROPOSAL
    
    Project Title: Machine Learning Pipeline for Real-time Document Classification
    Student: Alex Rivera
    Advisor: Dr. Robert Johnson
    
    Project Description:
    Develop an end-to-end machine learning system that can classify documents
    in real-time using deep learning techniques.
    
    Objectives:
    1. Design and train a neural network model
    2. Build REST API for inference
    3. Deploy on cloud infrastructure
    4. Achieve 95%+ accuracy on test dataset
    
    Timeline:
    January 2025: Research and dataset preparation
    February 2025: Model development and training
    March 2025: API development and testing
    April 2025: Deployment and evaluation
    May 2025: Final report and presentation
    
    Expected Deliverables:
    - Trained ML model
    - API documentation
    - Deployment guide
    - Performance report
    
    Submitted: December 29, 2024
    """),
    
    "Capstone_Data_Collection_Log.xlsx": ("CAPSTONE_WORK", """
    Data Collection Log for Capstone Project
    
    Date | Sample ID | Source | Features Collected | Status
    12/1 | 001 | Database A | 50 documents | Complete
    12/5 | 002 | Database B | 75 documents | Complete
    12/10 | 003 | Database C | 100 documents | Complete
    12/15 | 004 | Database D | 80 documents | Complete
    12/20 | 005 | Database E | 95 documents | Complete
    
    Total Samples Collected: 400 documents
    Data Cleaning Progress: 90% complete
    """),
    
    "Final_Presentation_Slides.pptx": ("CAPSTONE_WORK", """
    FINAL PRESENTATION SLIDES
    
    Slide 1: Title Slide
    Title: "Capstone Project Final Presentation"
    Student: Casey Thompson
    Date: December 29, 2024
    
    Slide 2: Project Overview
    - Problem statement
    - Objectives
    - Scope
    
    Slide 3: Methodology
    - Research approach
    - Implementation strategy
    - Tools and technologies
    
    Slide 4: Results
    - Key findings
    - Performance metrics
    - Comparative analysis
    
    Slide 5: Conclusion & Future Work
    - Summary of achievements
    - Recommendations
    - Next steps
    """),
}

print("\n📁 Creating dataset on Desktop...\n")

for filename, (category, content) in files_to_create.items():
    file_path = os.path.join(DESKTOP_PATH, filename)
    
    try:
        if filename.endswith(".pdf"):
            # Create PDF
            c = canvas.Canvas(file_path, pagesize=letter)
            c.setFont("Helvetica", 10)
            y = 750
            for line in content.split("\n"):
                if y < 50:
                    c.showPage()
                    y = 750
                c.drawString(50, y, line[:70])
                y -= 15
            c.save()
        
        elif filename.endswith(".docx"):
            # Create DOCX
            doc = Document()
            doc.add_paragraph(content)
            doc.save(file_path)
        
        elif filename.endswith(".xlsx"):
            # Create XLSX
            wb = Workbook()
            ws = wb.active
            for i, line in enumerate(content.split("\n"), 1):
                ws.append([line])
            wb.save(file_path)
        
        elif filename.endswith(".pptx"):
            # Create PPTX
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body_shape = slide.placeholders[1]
            title.text = filename.replace(".pptx", "")
            tf = body_shape.text_frame
            tf.text = content
            prs.save(file_path)
        
        else:
            # Create TXT or MD
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(content)
        
        print(f"✅ Created: {filename}")
    
    except Exception as e:
        print(f"❌ Failed to create {filename}: {e}")

print("\n" + "="*70)
print("✅ DATASET SETUP COMPLETE!")
print("="*70)
print(f"\n📁 Files created on Desktop:")
print(f"   Location: {DESKTOP_PATH}")
print(f"\n📄 FILES (10 total):")

files_on_desktop = [f for f in os.listdir(DESKTOP_PATH) 
                    if os.path.isfile(os.path.join(DESKTOP_PATH, f)) 
                    and f.endswith((".pdf", ".docx", ".xlsx", ".pptx", ".md", ".txt"))]
print(f"   Total: {len(files_on_desktop)} files\n")

for filename in sorted(files_on_desktop):
    print(f"   ✓ {filename}")

print(f"\n📂 FOLDERS (3 total - empty, ready for automation):")
for folder in FOLDERS:
    folder_path = os.path.join(DESKTOP_PATH, folder)
    if os.path.isdir(folder_path):
        file_count = len([f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))])
        print(f"   ✓ {folder} ({file_count} files inside)")

print("\n" + "="*70)
print("Ready for automation! All files are on Desktop root.")
print("Automation will classify and move them into the correct folders.")
print("="*70)
